from dataclasses import dataclass, field
from statistics import median
import logging
import mimetypes
import os
import re
import time
from docling.document_converter import DocumentConverter
import pandas as pd
import pytesseract
import requests
from docx import Document
from pdf2image import convert_from_path
from PIL import Image
from pypdf import PdfReader

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz
except ImportError:
    fitz = None

from pilotcore.config import TRACEPILOT_URL

logger = logging.getLogger(__name__)


class TextExtractionError(Exception):
    pass


@dataclass
class TextSection:
    text: str
    metadata: dict = field(default_factory=dict)


def detect_section_title(page_dict):
    """
    Return the highest-confidence heading detected on a page.
    Returns None if no heading meets the threshold.
    """

    candidates = []
    font_sizes = []

    # collect font sizes
    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue

        for line in block.get("lines", []):
            for span in line.get("spans", []):
                size = span.get("size")
                if size:
                    font_sizes.append(size)

    if not font_sizes:
        return None

    median_size = median(font_sizes)
    page_height = page_dict.get("height", 1)

    for block in page_dict.get("blocks", []):
        if block.get("type") != 0:
            continue

        lines = block.get("lines", [])
        if not lines:
            continue

        text_parts = []
        max_score = 0

        for line in lines:
            for span in line.get("spans", []):

                text = (span.get("text") or "").strip()

                if not text:
                    continue

                score = 0
                size = span.get("size", 0)
                flags = span.get("flags", 0)
                font = span.get("font", "")

                # font size
                ratio = size / median_size if median_size else 1

                if ratio >= 1.4:
                    score += 4
                elif ratio >= 1.1:
                    score += 2

                # bold flag
                if flags & 16:
                    score += 3

                # bold font name
                if any(token in font.lower() for token in ["bold", "medi"]):
                    score += 2

                # text length
                if len(text) <= 60:
                    score += 2
                elif len(text) <= 120:
                    score += 1
                else:
                    score -= 2

                # ALL CAPS
                if len(text) > 2 and text.isupper():
                    score += 2

                # Title Case
                if len(text.split()) > 1 and text.istitle():
                    score += 1

                text_parts.append(text)
                max_score = max(max_score, score)

        candidate_text = " ".join(text_parts).strip()

        if not candidate_text:
            continue

        if len(candidate_text.split()) > 12:
            continue

        # isolated block
        if len(lines) == 1:
            max_score += 1

        # top of page
        bbox = block.get("bbox", [0, 0, 0, 0])
        y0 = bbox[1]

        if y0 < page_height * 0.15:
            max_score += 1

        # reject lone numbers
        if re.fullmatch(r"[\d.]+", candidate_text):
            continue

        candidates.append(
            {
                "text": candidate_text,
                "score": max_score,
            }
        )

    if not candidates:
        return None

    candidates.sort(key=lambda x: x["score"], reverse=True)

    HEADING_WORD_LIMIT = 8
    filtered = [c for c in candidates if len(c["text"].split()) <= HEADING_WORD_LIMIT]

    best = filtered[0] if filtered else candidates[0]

    if best["score"] < 8:
        return None

    return best["text"]


def clean_text(text):
    text = (text or "").replace("\x00", "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(
    text,
    chunk_size=500,
    overlap=80,
):

    chunks = []

    start = 0

    while start < len(text):

        end = start + chunk_size

        if end < len(text):

            sentence_end = text.rfind(".", start, end)

            newline_end = text.rfind("\n", start, end)

            boundary = max(sentence_end, newline_end)

            if boundary != -1 and boundary > start:

                end = boundary + 1

        chunk = text[start:end].strip()

        if chunk:

            chunks.append(chunk)

        next_start = end - overlap

        # Ensure forward progress to avoid infinite loops
        if next_start <= start:

            next_start = end

        start = next_start

    return chunks


def detect_type(file_path, mime_type=None):
    extension = os.path.splitext(file_path)[1].lower()
    detected_mime = mime_type or mimetypes.guess_type(file_path)[0] or ""
    return extension, detected_mime


def extract_pdf_text(file_path):
    if fitz is not None:
        return extract_pdf_text_pymupdf(file_path)

    return extract_pdf_text_pypdf(file_path)


def extract_pdf_text_pymupdf(file_path):
    sections = []

    with fitz.open(file_path) as doc:
        for page_index, page in enumerate(doc, start=1):
            page_dict = page.get_text("dict")

            text = clean_text(page.get_text("text"))

            section_title = detect_section_title(page_dict)

            metadata = {
                "page": page_index,
            }

            if section_title:
                metadata["section_title"] = section_title

            if text:
                sections.append(
                    TextSection(
                        text=text,
                        metadata=metadata,
                    )
                )

        page_count = doc.page_count

    logger.info("PDF digital extraction processed %s pages with PyMuPDF", page_count)
    logger.info(
        "PDF digital extraction produced %s chars", sum(len(s.text) for s in sections)
    )
    return sections


def extract_pdf_text_pypdf(file_path):
    reader = PdfReader(file_path)
    sections = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        if text:
            sections.append(
                TextSection(
                    text=text,
                    metadata={"page": page_number},
                )
            )

    logger.info(
        "PDF digital extraction processed %s pages with pypdf", len(reader.pages)
    )
    logger.info(
        "PDF digital extraction produced %s chars", sum(len(s.text) for s in sections)
    )
    return sections


def extract_pdf_docling(file_path):
    try:
        converter = DocumentConverter()

        result = converter.convert(file_path)

        text = clean_text(result.document.export_to_markdown())

        if not text:
            return []

        return [
            TextSection(
                text=text,
                metadata={"extractor": "docling"},
            )
        ]

    except Exception as e:
        logger.exception(
            "Docling extraction failed: %s",
            e,
        )
        return []


def extract_pdf_ocr(file_path):
    try:
        images = convert_from_path(file_path)
        logger.info("OCR pages processed: %s", len(images))

        sections = []
        for page_number, image in enumerate(images, start=1):
            text = clean_text(pytesseract.image_to_string(image))
            if text:
                sections.append(
                    TextSection(
                        text=text,
                        metadata={"page": page_number, "ocr": True},
                    )
                )

        logger.info("OCR extracted %s chars", sum(len(s.text) for s in sections))
        return sections

    except Exception as e:
        logger.exception("OCR failed: %s", e)
        return []


def extract_pdf_sections(file_path):

    sections = extract_pdf_text(file_path)

    total_chars = sum(len(section.text) for section in sections)

    logger.info(
        "Digital PDF extraction produced %s chars",
        total_chars,
    )

    if total_chars > 500:
        return sections

    logger.info("PyMuPDF extraction weak, trying Docling")

    sections = extract_pdf_docling(file_path)

    docling_chars = sum(len(section.text) for section in sections)

    if docling_chars > 500:
        return sections

    logger.info("Docling failed, OCR triggered")

    return extract_pdf_ocr(file_path)


def extract_docx_sections(file_path):
    doc = Document(file_path)
    text = "\n".join(para.text for para in doc.paragraphs)

    return [
        TextSection(
            text=clean_text(text),
            metadata={},
        )
    ]


def extract_pptx_sections(file_path):
    if Presentation is None:
        raise TextExtractionError("PPTX extraction dependency is not installed")

    presentation = Presentation(file_path)
    sections = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    ]
                    if cells:
                        parts.append(" | ".join(cells))

        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(notes)
        except Exception:
            pass

        text = clean_text("\n".join(parts))
        if text:
            sections.append(
                TextSection(
                    text=text,
                    metadata={"slide": slide_number},
                )
            )

    logger.info("PPTX slides processed: %s", len(presentation.slides))
    return sections


def extract_txt_sections(file_path):
    with open(
        file_path,
        "r",
        encoding="utf-8",
        errors="ignore",
    ) as f:
        return [TextSection(text=clean_text(f.read()), metadata={})]


def extract_csv_sections(file_path):
    df = pd.read_csv(file_path)
    return dataframe_to_sections(df)


def extract_xlsx_sections(file_path):
    sheets = pd.read_excel(file_path, sheet_name=None)
    sections = []

    for sheet_name, df in sheets.items():
        for section in dataframe_to_sections(df):
            section.metadata["sheet"] = sheet_name
            sections.append(section)

    logger.info("XLSX sheets processed: %s", len(sheets))
    return sections


def dataframe_to_sections(df):
    sections = []
    df = df.fillna("")

    for row_number, row in df.iterrows():
        parts = []

        for column, value in row.items():
            value = str(value).strip()
            if value:
                parts.append(f"{column}: {value}")

        text = clean_text("\n".join(parts))
        if text:
            sections.append(
                TextSection(
                    text=text,
                    metadata={"row": int(row_number) + 1},
                )
            )

    return sections


def extract_image_sections(file_path):
    try:
        image = Image.open(file_path)
        text = clean_text(pytesseract.image_to_string(image))
        logger.info("Image OCR extracted %s chars", len(text))
        return [TextSection(text=text, metadata={"ocr": True})]

    except Exception as e:
        logger.exception("Image OCR failed: %s", e)
        return []


def extract_text_sections(file_path, mime_type=None):
    extension, detected_mime = detect_type(file_path, mime_type)

    extractors = {
        ".pdf": extract_pdf_sections,
        ".docx": extract_docx_sections,
        ".pptx": extract_pptx_sections,
        ".txt": extract_txt_sections,
        ".md": extract_txt_sections,
        ".csv": extract_csv_sections,
        ".xlsx": extract_xlsx_sections,
        ".png": extract_image_sections,
        ".jpg": extract_image_sections,
        ".jpeg": extract_image_sections,
        ".webp": extract_image_sections,
    }

    extractor = extractors.get(extension)

    if not extractor:
        raise TextExtractionError(
            f"Unsupported file type: {extension or detected_mime}"
        )

    logger.info(
        "Extractor selected: %s extension=%s mime=%s",
        extractor.__name__,
        extension,
        detected_mime,
    )

    try:
        sections = extractor(file_path)
    except TextExtractionError:
        raise
    except Exception as e:
        logger.exception("Extraction failed: %s", e)
        raise TextExtractionError("Could not extract text from document") from e

    cleaned_sections = []
    for section in sections:
        text = clean_text(section.text)
        if text:
            section.text = text
            cleaned_sections.append(section)

    extracted_length = sum(len(section.text) for section in cleaned_sections)
    logger.info("Extracted text length: %s", extracted_length)

    if not cleaned_sections:
        if extension == ".pdf":
            raise TextExtractionError("Could not extract text from PDF")
        raise TextExtractionError("Could not extract text from document")

    return cleaned_sections


def extract_text(file_path, mime_type=None):
    return "\n\n".join(
        section.text for section in extract_text_sections(file_path, mime_type)
    )


def process_document(
    file_path,
    user_id,
    document_id,
    mime_type=None,
):
    from DocPilot.backend.app.services.rag import add_chunks

    start_time = time.perf_counter()

    sections = extract_text_sections(file_path, mime_type)

    all_chunks = []
    source_file = os.path.basename(file_path)
    extension, _ = detect_type(file_path, mime_type)
    chunk_id = 0

    for section in sections:
        metadata = {
            "source_file": source_file,
            "file_type": extension.lstrip("."),
            **section.metadata,
        }

        page = metadata.get("page") or metadata.get("slide") or metadata.get("row") or 1

        for chunk in chunk_text(section.text):
            all_chunks.append(
                {
                    "document_id": document_id,
                    "text": chunk,
                    "source": source_file,
                    "source_file": source_file,
                    "file_type": extension.lstrip("."),
                    "page": page,
                    "chunk_id": chunk_id,
                    "metadata": metadata,
                }
            )
            chunk_id += 1

    logger.info("Chunk count: %s", len(all_chunks))

    if not all_chunks:
        raise TextExtractionError("Could not extract enough text from document")

    add_chunks(all_chunks, user_id)

    logger.info("Embedded %s chunks successfully", len(all_chunks))

    latency_ms = (time.perf_counter() - start_time) * 1000
    char_count = sum(len(section.text) for section in sections)

    try:
        resp = requests.post(
            f"{TRACEPILOT_URL}/tracepilot/ingest/document",
            json={
                "document_id": str(document_id),
                "user_id": str(user_id),
                "filename": source_file,
                "chunk_count": len(all_chunks),
                "char_count": char_count,
                "latency_ms": round(latency_ms, 2),
                "status": "success",
            },
            timeout=5,
        )

        logger.info("[TracePilot] document ingest status=%s", resp.status_code)
        logger.info("[TracePilot] response=%s", resp.text)

    except Exception as e:
        logger.info("[TracePilot] document ingest failed: %r", e)
